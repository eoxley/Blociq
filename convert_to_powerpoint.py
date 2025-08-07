#!/usr/bin/env python3
"""
Convert BlocIQ HTML Pitch Deck to PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json

def create_blociq_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define BlocIQ brand colors
    BLOCIQ_TEAL = RGBColor(0, 140, 143)  # #008C8F
    BLOCIQ_PURPLE = RGBColor(118, 69, 237)  # #7645ED
    BLOCIQ_GRADIENT_START = RGBColor(0, 140, 143)
    BLOCIQ_GRADIENT_END = RGBColor(118, 69, 237)
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add logo placeholder
    logo_left = Inches(1)
    logo_top = Inches(0.5)
    logo_width = Inches(2)
    logo_height = Inches(1)
    
    logo_shape = slide.shapes.add_textbox(logo_left, logo_top, logo_width, logo_height)
    logo_frame = logo_shape.text_frame
    logo_frame.text = "BlocIQ"
    logo_para = logo_frame.paragraphs[0]
    logo_para.font.size = Pt(36)
    logo_para.font.bold = True
    logo_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add title
    title_left = Inches(1)
    title_top = Inches(2)
    title_width = Inches(8)
    title_height = Inches(2)
    
    title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_shape.text_frame
    title_frame.text = "From Inbox Chaos to Intelligent Block Management"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add subtitle
    subtitle_left = Inches(1)
    subtitle_top = Inches(4)
    subtitle_width = Inches(8)
    subtitle_height = Inches(1)
    
    subtitle_shape = slide.shapes.add_textbox(subtitle_left, subtitle_top, subtitle_width, subtitle_height)
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "AI-powered compliance for modern property managers"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = BLOCIQ_PURPLE
    
    # Add content
    content_left = Inches(1)
    content_top = Inches(5.5)
    content_width = Inches(8)
    content_height = Inches(3)
    
    content_shape = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
    content_frame = content_shape.text_frame
    content_frame.text = "While AI is being rapidly adopted across industries, most property firms are using generic tools like ChatGPT without oversight, data safeguards, or regulatory alignment. That's a risk — for firms, for clients, and for residents.\n\nWe're building the operating system for modern property teams — from prime city blocks to social housing estates."
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(14)
    
    # Slide 2: The Problem
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(title_left, title_top - Inches(1), title_width, title_height)
    title_frame = title_shape.text_frame
    title_frame.text = "The Problem"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add subtitle
    subtitle_shape = slide.shapes.add_textbox(subtitle_left, subtitle_top - Inches(1), subtitle_width, subtitle_height)
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "Property management is drowning in chaos"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = BLOCIQ_PURPLE
    
    # Add problem points
    problems = [
        "Inbox Overload: Property managers spend 40% of their day sorting through emails",
        "Compliance Nightmare: Manual tracking of hundreds of compliance items across multiple buildings",
        "Generic AI Risk: Teams using ChatGPT for sensitive property data without proper safeguards"
    ]
    
    for i, problem in enumerate(problems):
        problem_left = Inches(1)
        problem_top = Inches(3 + i * 1.2)
        problem_width = Inches(8)
        problem_height = Inches(1)
        
        problem_shape = slide.shapes.add_textbox(problem_left, problem_top, problem_width, problem_height)
        problem_frame = problem_shape.text_frame
        problem_frame.text = f"• {problem}"
        problem_para = problem_frame.paragraphs[0]
        problem_para.font.size = Pt(16)
        problem_para.font.bold = True
        problem_para.font.color.rgb = BLOCIQ_TEAL
    
    # Slide 3: The Solution
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(title_left, title_top - Inches(1), title_width, title_height)
    title_frame = title_shape.text_frame
    title_frame.text = "The Solution"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add subtitle
    subtitle_shape = slide.shapes.add_textbox(subtitle_left, subtitle_top - Inches(1), subtitle_width, subtitle_height)
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "BlocIQ: AI-powered compliance for modern property managers"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = BLOCIQ_PURPLE
    
    # Add solution features
    features = [
        ("Intelligent Inbox", "AI-powered email sorting, prioritization, and automated responses"),
        ("Compliance Tracking", "Automated compliance monitoring with deadline alerts and risk assessment"),
        ("AI Assistant", "Context-aware AI that understands property management and compliance"),
        ("GDPR Compliant", "Built-in data protection and regulatory compliance safeguards")
    ]
    
    for i, (feature, description) in enumerate(features):
        feature_left = Inches(1 + (i % 2) * 4.5)
        feature_top = Inches(3 + (i // 2) * 1.5)
        feature_width = Inches(4)
        feature_height = Inches(1.2)
        
        feature_shape = slide.shapes.add_textbox(feature_left, feature_top, feature_width, feature_height)
        feature_frame = feature_shape.text_frame
        feature_frame.text = f"{feature}\n{description}"
        feature_para = feature_frame.paragraphs[0]
        feature_para.font.size = Pt(14)
        feature_para.font.bold = True
        feature_para.font.color.rgb = BLOCIQ_TEAL
        
        desc_para = feature_frame.paragraphs[1]
        desc_para.font.size = Pt(12)
    
    # Slide 4: Market Opportunity
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(title_left, title_top - Inches(1), title_width, title_height)
    title_frame = title_shape.text_frame
    title_frame.text = "Market Opportunity"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add subtitle
    subtitle_shape = slide.shapes.add_textbox(subtitle_left, subtitle_top - Inches(1), subtitle_width, subtitle_height)
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "A £2.3B market ready for disruption"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = BLOCIQ_PURPLE
    
    # Add market stats
    stats = [
        ("£2.3B", "Market Size"),
        ("4,500", "Property Firms"),
        ("2.9M", "Leasehold Properties")
    ]
    
    for i, (stat, label) in enumerate(stats):
        stat_left = Inches(1 + i * 2.5)
        stat_top = Inches(3)
        stat_width = Inches(2)
        stat_height = Inches(1.5)
        
        stat_shape = slide.shapes.add_textbox(stat_left, stat_top, stat_width, stat_height)
        stat_frame = stat_shape.text_frame
        stat_frame.text = f"{stat}\n{label}"
        stat_para = stat_frame.paragraphs[0]
        stat_para.font.size = Pt(24)
        stat_para.font.bold = True
        stat_para.font.color.rgb = BLOCIQ_TEAL
        stat_para.alignment = PP_ALIGN.CENTER
        
        label_para = stat_frame.paragraphs[1]
        label_para.font.size = Pt(14)
        label_para.alignment = PP_ALIGN.CENTER
    
    # Slide 5: Technology Stack
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(title_left, title_top - Inches(1), title_width, title_height)
    title_frame = title_shape.text_frame
    title_frame.text = "Technology & Compliance Stack"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add subtitle
    subtitle_shape = slide.shapes.add_textbox(subtitle_left, subtitle_top - Inches(1), subtitle_width, subtitle_height)
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "Built for security, compliance, and scale"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = BLOCIQ_PURPLE
    
    # Add tech stack
    tech_stack = [
        ("Supabase", "Secure, GDPR-compliant database with real-time capabilities"),
        ("OpenAI GPT-4", "Advanced AI with property management context"),
        ("Microsoft Graph", "Enterprise-grade email integration"),
        ("Next.js 14", "Modern, performant web framework")
    ]
    
    for i, (tech, description) in enumerate(tech_stack):
        tech_left = Inches(1 + (i % 2) * 4.5)
        tech_top = Inches(3 + (i // 2) * 1.5)
        tech_width = Inches(4)
        tech_height = Inches(1.2)
        
        tech_shape = slide.shapes.add_textbox(tech_left, tech_top, tech_width, tech_height)
        tech_frame = tech_shape.text_frame
        tech_frame.text = f"{tech}\n{description}"
        tech_para = tech_frame.paragraphs[0]
        tech_para.font.size = Pt(14)
        tech_para.font.bold = True
        tech_para.font.color.rgb = BLOCIQ_TEAL
        
        desc_para = tech_frame.paragraphs[1]
        desc_para.font.size = Pt(12)
    
    # Slide 6: Revenue Model
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(title_left, title_top - Inches(1), title_width, title_height)
    title_frame = title_shape.text_frame
    title_frame.text = "Revenue Model"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add subtitle
    subtitle_shape = slide.shapes.add_textbox(subtitle_left, subtitle_top - Inches(1), subtitle_width, subtitle_height)
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "SaaS subscription with clear value proposition"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = BLOCIQ_PURPLE
    
    # Add pricing tiers
    pricing = [
        ("Starter", "£99/month", "Up to 5 buildings"),
        ("Professional", "£299/month", "Up to 25 buildings"),
        ("Enterprise", "Custom", "Unlimited buildings")
    ]
    
    for i, (tier, price, description) in enumerate(pricing):
        tier_left = Inches(1 + i * 2.5)
        tier_top = Inches(3)
        tier_width = Inches(2)
        tier_height = Inches(1.5)
        
        tier_shape = slide.shapes.add_textbox(tier_left, tier_top, tier_width, tier_height)
        tier_frame = tier_shape.text_frame
        tier_frame.text = f"{tier}\n{price}\n{description}"
        tier_para = tier_frame.paragraphs[0]
        tier_para.font.size = Pt(16)
        tier_para.font.bold = True
        tier_para.font.color.rgb = BLOCIQ_TEAL
        tier_para.alignment = PP_ALIGN.CENTER
        
        price_para = tier_frame.paragraphs[1]
        price_para.font.size = Pt(20)
        price_para.font.bold = True
        price_para.font.color.rgb = BLOCIQ_TEAL
        price_para.alignment = PP_ALIGN.CENTER
        
        desc_para = tier_frame.paragraphs[2]
        desc_para.font.size = Pt(12)
        desc_para.alignment = PP_ALIGN.CENTER
    
    # Slide 7: Traction
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(title_left, title_top - Inches(1), title_width, title_height)
    title_frame = title_shape.text_frame
    title_frame.text = "Traction"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add subtitle
    subtitle_shape = slide.shapes.add_textbox(subtitle_left, subtitle_top - Inches(1), subtitle_width, subtitle_height)
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "Early validation and growth momentum"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = BLOCIQ_PURPLE
    
    # Add traction metrics
    metrics = [
        ("15", "Pilot customers"),
        ("£25K", "ARR"),
        ("95%", "Retention rate")
    ]
    
    for i, (metric, label) in enumerate(metrics):
        metric_left = Inches(1 + i * 2.5)
        metric_top = Inches(3)
        metric_width = Inches(2)
        metric_height = Inches(1.5)
        
        metric_shape = slide.shapes.add_textbox(metric_left, metric_top, metric_width, metric_height)
        metric_frame = metric_shape.text_frame
        metric_frame.text = f"{metric}\n{label}"
        metric_para = metric_frame.paragraphs[0]
        metric_para.font.size = Pt(24)
        metric_para.font.bold = True
        metric_para.font.color.rgb = BLOCIQ_TEAL
        metric_para.alignment = PP_ALIGN.CENTER
        
        label_para = metric_frame.paragraphs[1]
        label_para.font.size = Pt(14)
        label_para.alignment = PP_ALIGN.CENTER
    
    # Slide 8: Founder Story
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(title_left, title_top - Inches(1), title_width, title_height)
    title_frame = title_shape.text_frame
    title_frame.text = "Founder Story"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add subtitle
    subtitle_shape = slide.shapes.add_textbox(subtitle_left, subtitle_top - Inches(1), subtitle_width, subtitle_height)
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "Ellie Oxley - Building the future of property management"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = BLOCIQ_PURPLE
    
    # Add founder story
    story_left = Inches(1)
    story_top = Inches(3)
    story_width = Inches(8)
    story_height = Inches(4)
    
    story_shape = slide.shapes.add_textbox(story_left, story_top, story_width, story_height)
    story_frame = story_shape.text_frame
    story_frame.text = "After 8 years in property management, Ellie experienced firsthand the chaos of managing hundreds of compliance items across multiple buildings. The industry was drowning in emails, missing deadlines, and using generic AI tools that didn't understand property management.\n\nBlocIQ was born from frustration with the status quo and a vision for intelligent, compliant property management.\n\n'I've lived the problem. Now I'm building the solution.'"
    story_para = story_frame.paragraphs[0]
    story_para.font.size = Pt(14)
    
    # Slide 9: Why Now / Why Pi Labs
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(title_left, title_top - Inches(1), title_width, title_height)
    title_frame = title_shape.text_frame
    title_frame.text = "Why Now & Why Pi Labs"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add subtitle
    subtitle_shape = slide.shapes.add_textbox(subtitle_left, subtitle_top - Inches(1), subtitle_width, subtitle_height)
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = "Perfect timing for the right partnership"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(18)
    subtitle_para.font.color.rgb = BLOCIQ_PURPLE
    
    # Add why now points
    why_now_points = [
        "AI Adoption: Property firms ready for AI solutions",
        "Regulatory Pressure: Increasing compliance requirements",
        "Cost Pressure: Need for operational efficiency",
        "Staff Shortages: Technology to augment teams"
    ]
    
    for i, point in enumerate(why_now_points):
        point_left = Inches(1)
        point_top = Inches(3 + i * 0.8)
        point_width = Inches(8)
        point_height = Inches(0.6)
        
        point_shape = slide.shapes.add_textbox(point_left, point_top, point_width, point_height)
        point_frame = point_shape.text_frame
        point_frame.text = f"• {point}"
        point_para = point_frame.paragraphs[0]
        point_para.font.size = Pt(14)
        point_para.font.bold = True
        point_para.font.color.rgb = BLOCIQ_TEAL
    
    # Add call to action
    cta_left = Inches(3)
    cta_top = Inches(6.5)
    cta_width = Inches(4)
    cta_height = Inches(1)
    
    cta_shape = slide.shapes.add_textbox(cta_left, cta_top, cta_width, cta_height)
    cta_frame = cta_shape.text_frame
    cta_frame.text = "Let's Build the Future Together"
    cta_para = cta_frame.paragraphs[0]
    cta_para.font.size = Pt(18)
    cta_para.font.bold = True
    cta_para.font.color.rgb = BLOCIQ_TEAL
    cta_para.alignment = PP_ALIGN.CENTER
    
    return prs

if __name__ == "__main__":
    # Create the presentation
    presentation = create_blociq_presentation()
    
    # Save the presentation
    output_file = "BlocIQ_Pitch_Deck.pptx"
    presentation.save(output_file)
    
    print(f"✅ PowerPoint presentation created: {output_file}")
    print("🎯 The presentation includes all 9 slides with BlocIQ branding")
    print("📊 Ready for investor presentations and demos") 